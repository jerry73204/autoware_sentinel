#!/usr/bin/env python3
"""Generate nano-ros-safety-features.pptx from template.pptx.

Uses "Simple Light" theme layouts (second set):
  - Layout 11 (TITLE) for the title slide
  - Layout 12 (TITLE_AND_BODY) for content slides

Features:
  - Syntax-highlighted code blocks (Pygments)
  - Native PowerPoint tables
  - Diagram images generated via Pillow
  - Consistent 10pt Courier New for code, 14pt body, 24pt titles
"""

from pathlib import Path
from io import BytesIO

from lxml import etree
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

from pygments import lex
from pygments.lexers import RustLexer, CLexer, BashLexer, TOMLLexer
from pygments.token import Token

HERE = Path(__file__).parent
TEMPLATE = HERE / "template.pptx"
OUTPUT = HERE / "nano-ros-safety-features.pptx"

# Layout indices — "Simple Light" second set
LAYOUT_TITLE = 11
LAYOUT_TITLE_BODY = 12

# Typography
FONT_CODE = "Courier New"
CODE_SIZE = Pt(10)
TITLE_SIZE = Pt(24)
BODY_SIZE = Pt(14)
SMALL_SIZE = Pt(12)
BULLET_SIZE = Pt(13)

# Pygments token → RGB color map (VS Code–inspired light theme)
TOKEN_COLORS = {
    Token.Keyword: RGBColor(0xAF, 0x00, 0xDB),          # purple
    Token.Keyword.Type: RGBColor(0x26, 0x7F, 0x99),      # teal
    Token.Name.Function: RGBColor(0x79, 0x5E, 0x26),     # dark yellow
    Token.Name.Decorator: RGBColor(0x79, 0x5E, 0x26),
    Token.Name.Attribute: RGBColor(0x00, 0x70, 0xC1),    # blue
    Token.Name.Builtin: RGBColor(0x26, 0x7F, 0x99),
    Token.String: RGBColor(0xA3, 0x11, 0x15),            # red
    Token.Literal.String: RGBColor(0xA3, 0x11, 0x15),
    Token.Number: RGBColor(0x09, 0x86, 0x58),            # green
    Token.Literal.Number: RGBColor(0x09, 0x86, 0x58),
    Token.Comment: RGBColor(0x6A, 0x99, 0x55),           # green
    Token.Comment.Single: RGBColor(0x6A, 0x99, 0x55),
    Token.Comment.Multiline: RGBColor(0x6A, 0x99, 0x55),
    Token.Operator: RGBColor(0x00, 0x00, 0x00),
    Token.Punctuation: RGBColor(0x00, 0x00, 0x00),
    Token.Name: RGBColor(0x00, 0x70, 0xC1),              # blue for names
}
DEFAULT_CODE_COLOR = RGBColor(0x1A, 0x1A, 0x2E)

LEXERS = {
    "rust": RustLexer(),
    "c": CLexer(),
    "bash": BashLexer(),
    "toml": TOMLLexer(),
    "text": None,  # no highlighting
}


# ============================================================================
# Helpers
# ============================================================================

def _suppress_bullet(para):
    """Remove bullet/numbering from a paragraph."""
    pPr = para._p.get_or_add_pPr()
    for tag in ["buNone", "buChar", "buAutoNum", "buFont", "buSzPct", "buSzPts"]:
        for el in pPr.findall(qn(f"a:{tag}")):
            pPr.remove(el)
    pPr.append(etree.SubElement(pPr, qn("a:buNone")))
    pPr.set("indent", "0")
    pPr.set("marL", "0")


def _get_token_color(ttype):
    """Walk up the token type hierarchy to find a matching color."""
    while ttype:
        if ttype in TOKEN_COLORS:
            return TOKEN_COLORS[ttype]
        ttype = ttype.parent
    return DEFAULT_CODE_COLOR


def _add_run(para, text, font_name=None, font_size=None, bold=False, color=None, italic=False):
    """Add a formatted run to a paragraph."""
    run = para.add_run()
    run.text = text
    if font_name:
        run.font.name = font_name
    if font_size:
        run.font.size = font_size
    if bold:
        run.font.bold = True
    if italic:
        run.font.italic = True
    if color:
        run.font.color.rgb = color
    return run


def add_body_para(tf, text, size=BODY_SIZE, bold=False, space_before=Pt(4), space_after=Pt(2)):
    """Add a bullet-point body text paragraph."""
    p = tf.add_paragraph()
    p.level = 0
    p.space_before = space_before
    p.space_after = space_after
    _add_run(p, text, font_size=size, bold=bold)
    return p


def add_bullet(tf, text, size=BULLET_SIZE, bold_prefix=None):
    """Add a bullet point. If bold_prefix given, first part is bold."""
    p = tf.add_paragraph()
    p.level = 0
    p.space_before = Pt(2)
    p.space_after = Pt(2)
    if bold_prefix:
        _add_run(p, bold_prefix, font_size=size, bold=True)
        _add_run(p, text, font_size=size)
    else:
        _add_run(p, text, font_size=size)
    return p


def add_code_block(tf, code, lang="rust", size=CODE_SIZE, space_before=Pt(6)):
    """Add a syntax-highlighted code block (no bullets)."""
    lexer = LEXERS.get(lang)
    for i, line in enumerate(code.strip("\n").split("\n")):
        p = tf.add_paragraph()
        p.space_before = Pt(0) if i > 0 else space_before
        p.space_after = Pt(0)
        p.level = 0
        _suppress_bullet(p)

        if lexer is None or not line.strip():
            # Plain text or empty line
            _add_run(p, line or " ", font_name=FONT_CODE, font_size=size, color=DEFAULT_CODE_COLOR)
        else:
            tokens = list(lex(line + "\n", lexer))
            for ttype, value in tokens:
                value = value.rstrip("\n")
                if not value:
                    continue
                color = _get_token_color(ttype)
                _add_run(p, value, font_name=FONT_CODE, font_size=size, color=color)


def add_table(slide, left, top, width, headers, rows, font_size=Pt(10)):
    """Add a native PowerPoint table to the slide."""
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.3 * n_rows))
    table = table_shape.table

    # Set column widths proportionally
    total_chars = sum(
        max(len(h), max((len(str(r[ci])) for r in rows), default=0))
        for ci, h in enumerate(headers)
    )
    for ci, h in enumerate(headers):
        col_chars = max(len(h), max((len(str(r[ci])) for r in rows), default=0))
        table.columns[ci].width = int(width * col_chars / total_chars)

    def _set_cell(cell, text, bold=False):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = font_size
        p.font.name = "Calibri"
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.name = "Calibri"
        if bold:
            run.font.bold = True
        # Minimal padding
        cell.margin_top = Pt(2)
        cell.margin_bottom = Pt(2)
        cell.margin_left = Pt(4)
        cell.margin_right = Pt(4)

    # Header row
    for ci, h in enumerate(headers):
        _set_cell(table.cell(0, ci), h, bold=True)

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            _set_cell(table.cell(ri + 1, ci), str(val))

    return table_shape


def _make_stack_diagram(layers, width=800, row_height=52, padding=16, font_size=18):
    """Create a defense-in-depth stack diagram image using Pillow."""
    height = len(layers) * row_height + padding * 2
    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)

    # Try to load a nice font, fall back to default
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", font_size)
        font_small = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", font_size - 4)
    except OSError:
        font = ImageFont.load_default()
        font_small = font

    # Color palette (top = darkest)
    colors = [
        "#2D3748", "#3182CE", "#2B6CB0", "#2C7A7B",
        "#276749", "#744210", "#9B2C2C",
    ]

    for i, (title, subtitle) in enumerate(layers):
        y = padding + i * row_height
        x = padding
        box_w = width - 2 * padding
        box_h = row_height - 4

        color = colors[i % len(colors)]
        draw.rounded_rectangle([x, y, x + box_w, y + box_h], radius=6, fill=color)

        # Center text
        title_bbox = draw.textbbox((0, 0), title, font=font)
        tw = title_bbox[2] - title_bbox[0]
        sub_bbox = draw.textbbox((0, 0), subtitle, font=font_small)
        sw = sub_bbox[2] - sub_bbox[0]

        cx = x + box_w // 2
        if subtitle:
            draw.text((cx - tw // 2, y + 4), title, fill="white", font=font)
            draw.text((cx - sw // 2, y + 4 + font_size + 1), subtitle, fill="#CBD5E0", font=font_small)
        else:
            draw.text((cx - tw // 2, y + (box_h - font_size) // 2), title, fill="white", font=font)

    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def new_slide(prs, layout_idx, title_text):
    """Create a new slide with title set and font size forced to TITLE_SIZE."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    ph = slide.placeholders[0]
    ph.text = title_text
    # Force title font size
    for para in ph.text_frame.paragraphs:
        for run in para.runs:
            if layout_idx != LAYOUT_TITLE:
                run.font.size = TITLE_SIZE
    return slide


def get_body(slide):
    """Get the body text frame (placeholder index 1)."""
    return slide.placeholders[1].text_frame


def clear_body(tf):
    """Clear default paragraphs from text frame."""
    for p in tf.paragraphs:
        p.text = ""


# ============================================================================
# Build slides
# ============================================================================

def build():
    prs = Presentation(str(TEMPLATE))

    # Delete existing template slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # =================================================================
    # Title Slide
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE, "nano-ros Safety Features")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "A Multi-Layered Approach to Safety-Critical Robotics", size=Pt(16))
    add_body_para(tf, "Hsiang-Jui (Jerry) Lin, National Taiwan University", size=Pt(12))

    # =================================================================
    # Safety at Every Layer (with native table)
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Safety at Every Layer")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "nano-ros applies defense-in-depth across the full stack.", size=BODY_SIZE)

    add_table(
        slide,
        left=Inches(0.5), top=Inches(2.0), width=Inches(9.0),
        headers=["Layer", "Mechanism", "Guarantee"],
        rows=[
            ["Language", "Rust ownership + borrowing", "Memory safety without GC"],
            ["Concurrency", "async/await \u2192 state machines", "Zero-overhead, known stack size"],
            ["Compiler", "#![no_std], no alloc, feature gates", "No heap, no OS dependency"],
            ["Types", "heapless::Vec<T,N>, FixedSequence", "Bounded memory at compile time"],
            ["Serialization", "CDR with bounds-checked writes", "No buffer overflows"],
            ["Scheduling", "Deterministic executor, RTIC", "Predictable timing"],
            ["Networking", "E2E CRC-32 + seq (EN 50159)", "Corruption/loss detection"],
            ["TSN", "ThreadX + NetX Duo driver", "Hardware-enforced determinism"],
            ["Verification", "67 Verus + 82 Kani proofs", "Mathematical correctness"],
        ],
    )

    # =================================================================
    # Rust Concurrency Without Overhead
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Rust Concurrency Without Overhead")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "Rust\u2019s ownership system enforces thread safety at compile time \u2014 no runtime cost.")
    add_code_block(tf, """\
// Ownership prevents data races: only one mutable reference at a time
let mut msg = Odometry::default();
let ref1 = &mut msg;
// let ref2 = &mut msg;  // Compile error: cannot borrow twice""", lang="rust")
    add_code_block(tf, """\
// Send + Sync traits — compiler checks thread-safety of every type
fn spawn_task<F: Send + 'static>(f: F) { /* ... */ }

spawn_task(move || {
    publisher.publish(&msg);  // Publisher must be Send — compiler verifies
});""", lang="rust")
    add_body_para(tf, "In C/C++, data races are runtime bugs. In Rust, they are compile errors.", space_before=Pt(8))
    add_body_para(tf, "No mutexes for single-owner data. No thread sanitizer \u2014 the compiler is the sanitizer.", size=SMALL_SIZE)

    # =================================================================
    # Async Without a Runtime (page 1)
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Async Without a Runtime")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "Rust provides async/await syntax but no built-in runtime. The language gives you the abstraction; you choose the executor.")
    add_code_block(tf, """\
// This reads like threaded code...
async fn service_call(client: &mut Client<OperateMrm>)
    -> Result<Response, NodeError>
{
    let promise = client.call(&request)?;
    promise.await   // Suspends here, resumes when reply arrives
}

// ...but compiles to a state machine (no thread, no stack allocation):
enum ServiceCallFuture {
    WaitingForReply { client: ..., promise: ... },
    Done,
}""", lang="rust")

    # =================================================================
    # Async Without a Runtime (page 2) — Poll + Waker
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Async Without a Runtime \u2014 Poll + Waker")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "The key insight: Promise implements Future using poll + waker.")
    add_code_block(tf, """\
impl<T> core::future::Future for Promise<'_, T> {
    type Output = Result<T, NodeError>;

    fn poll(self: Pin<&mut Self>, cx: &mut Context<'_>)
        -> Poll<Self::Output>
    {
        match self.get_mut().try_recv() {
            Ok(Some(reply)) => Poll::Ready(Ok(reply)),
            Ok(None) => {
                // Wake me when data arrives
                self.handle.register_waker(cx.waker());
                Poll::Pending
            }
            Err(e) => Poll::Ready(Err(e)),
        }
    }
}""", lang="rust")
    add_body_para(tf, "No busy-waiting. The waker notifies the executor only when data is ready.", size=SMALL_SIZE, space_before=Pt(8))

    # =================================================================
    # Why Async Matters for Safety (page 1)
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Why Async Matters for Safety")
    tf = get_body(slide)
    clear_body(tf)
    add_bullet(tf, "Stack size is known at compile time. ", bold_prefix="Known size: ", size=BODY_SIZE)
    add_body_para(tf, "Each async fn compiles to a state machine enum \u2014 the compiler computes its exact size. No per-task stack allocation.", size=SMALL_SIZE)
    add_bullet(tf, "Multiple \u201cconcurrent\u201d tasks share one thread:", bold_prefix="Runs on a single-threaded MCU. ", size=BODY_SIZE)
    add_code_block(tf, """\
// nano-ros executor: drives all async tasks on a single thread
pub async fn spin_async(&mut self) -> ! {
    loop {
        self.spin_once(1);   // Process all ready callbacks
        core::future::poll_fn::<(), _>(|cx| {
            cx.waker().wake_by_ref();
            Poll::Pending      // Yield to other tasks
        }).await;
    }
}""", lang="rust")

    # =================================================================
    # Why Async Matters for Safety (page 2)
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Why Async Matters for Safety \u2014 No select()")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "Eliminates error-prone select()/poll() with fd sets. Rust async composes naturally:")
    add_code_block(tf, """\
// RTIC: two independent async tasks, zero shared state, zero locks
#[task(local = [executor], priority = 1)]
async fn net_poll(cx: net_poll::Context) {
    loop {
        cx.local.executor.spin_once(0);
        Mono::delay(10.millis()).await;
    }
}

#[task(local = [publisher], priority = 1)]
async fn publish(cx: publish::Context) {
    loop {
        cx.local.publisher.publish(&Int32 { data: 42 }).ok();
        Mono::delay(1000.millis()).await;
    }
}""", lang="rust")
    add_body_para(tf, "The compiler generates the same code as hand-written state machines, but with the readability of sequential code.", size=SMALL_SIZE, space_before=Pt(8))

    # =================================================================
    # Compiler-Enforced no_std and Alloc-Free
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Compiler-Enforced no_std and Alloc-Free")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "Every core crate compiles without the standard library. std and alloc are opt-in.")
    add_code_block(tf, """\
# nros-core/Cargo.toml
[features]
default = ["std"]
std = ["nros-serdes/std"]
alloc = ["nros-serdes/alloc"]""", lang="toml")
    add_body_para(tf, "Algorithm crates must cross-compile to bare-metal ARM:", space_before=Pt(8))
    add_code_block(tf, "cargo check --target thumbv7em-none-eabihf", lang="bash")
    add_body_para(tf, "If any code pulls in std or alloc, the compiler rejects it.", bold=True, space_before=Pt(8))
    add_body_para(tf, "This is not a convention \u2014 it is a hard gate enforced on every CI run.", size=SMALL_SIZE)

    # =================================================================
    # Compiler-Enforced Platform Selection
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Compiler-Enforced Platform Selection")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "Platform backends are mutually exclusive feature flags, not runtime configuration.")
    add_code_block(tf, """\
# nros-node/Cargo.toml — select exactly one
platform-posix      = [...]   # Linux, macOS
platform-zephyr     = [...]   # Zephyr RTOS
platform-freertos   = [...]   # FreeRTOS
platform-nuttx      = [...]   # NuttX
platform-threadx    = [...]   # Azure RTOS / ThreadX
platform-bare-metal = [...]   # No OS""", lang="toml")
    add_body_para(tf, "Same principle for RMW (middleware) backends:", space_before=Pt(8))
    add_code_block(tf, """\
rmw-zenoh = ["dep:nros-rmw-zenoh"]   # zenoh-pico (primary)
rmw-xrce  = ["dep:nros-rmw-xrce"]   # Micro-XRCE-DDS
rmw-dds   = ["dep:nros-rmw-dds"]     # Direct DDS
rmw-cffi  = ["dep:nros-rmw-cffi"]    # C FFI bridge""", lang="toml")
    add_body_para(tf, "Wrong platform + wrong RMW = link error at build time, not runtime crash.", bold=True, space_before=Pt(8))

    # =================================================================
    # Safe Types in Generated Message Code (Rust) — page 1: struct
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Safe Types in Generated Message Code")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "ROS 2 messages are code-generated into #![no_std] Rust crates. Dynamic containers are replaced with fixed-capacity alternatives.")
    add_code_block(tf, """\
// Generated from rcl_interfaces/srv/SetParameters
pub struct SetParametersRequest {
    pub parameters: heapless::Vec<Parameter, 64>,
}""", lang="rust")
    add_bullet(tf, " heapless::Vec<T, 64> holds at most 64 elements \u2014 stack-allocated", bold_prefix="Bounded:")
    add_bullet(tf, " CapacityExceeded error, not silent truncation or panic", bold_prefix="Explicit overflow:")
    add_bullet(tf, " Works on bare-metal MCUs with zero heap", bold_prefix="No allocator:")

    # =================================================================
    # Safe Types — page 2: Deserialize impl
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Safe Types \u2014 Deserialize with Capacity Check")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "The generated Deserialize impl checks capacity on every push:")
    add_code_block(tf, """\
impl Deserialize for SetParametersRequest {
    fn deserialize(reader: &mut CdrReader) -> Result<Self, DeserError> {
        let len = reader.read_u32()? as usize;
        let mut vec = heapless::Vec::new();
        for _ in 0..len {
            vec.push(Deserialize::deserialize(reader)?)
                .map_err(|_| DeserError::CapacityExceeded)?;
        }
        Ok(Self { parameters: vec })
    }
}""", lang="rust")
    add_body_para(tf, "Overflow returns CapacityExceeded \u2014 not a panic, not silent truncation.", size=SMALL_SIZE, space_before=Pt(8))

    # =================================================================
    # C++ Safe Types — page 1: FixedSequence
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "C++ Safe Types \u2014 Freestanding C++14")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "The same fixed-capacity pattern applies to the C++ API. No STL, no exceptions, no RTTI.")
    add_code_block(tf, """\
// nros::FixedSequence<T, N> — C++ equivalent of heapless::Vec<T, N>
template <typename T, size_t N> struct FixedSequence {
    uint32_t size;
    T data[N];

    bool push_back(const T& val) {
        if (size >= N) return false;  // Capacity check — no heap
        data[size++] = val;
        return true;
    }
};""", lang="c")
    add_body_para(tf, "Memory layout is repr(C) \u2014 identical on Rust and C++ sides, direct FFI.", size=SMALL_SIZE, space_before=Pt(6))

    # =================================================================
    # C++ Safe Types — page 2: NROS_TRY error handling
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "C++ Safe Types \u2014 Error Handling")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "NROS_TRY macro replaces try/catch \u2014 propagates errors via return codes:")
    add_code_block(tf, """\
nros::Result run() {
    nros::Node node;
    NROS_TRY(nros::create_node(node, "my_node"));

    nros::Publisher<std_msgs::msg::Int32> pub;
    NROS_TRY(node.create_publisher(pub, "/chatter"));

    std_msgs::msg::Int32 msg;
    msg.data = 42;
    NROS_TRY(pub.publish(msg));
    return nros::Result::success();
}""", lang="c")
    add_body_para(tf, "No exceptions, no RTTI \u2014 works on freestanding C++14 targets.", size=SMALL_SIZE, space_before=Pt(6))

    # =================================================================
    # Bounds-Checked CDR Serialization — page 1: struct
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Bounds-Checked CDR Serialization")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "The CDR encoder checks every write against the buffer.")
    add_code_block(tf, """\
pub struct CdrWriter<'a> {
    buf: &'a mut [u8],   // Fixed buffer — no reallocation
    pos: usize,
    origin: usize,       // Payload origin for alignment
}""", lang="rust")
    add_body_para(tf, "Fixed-size buffer \u2014 no Vec, no reallocation, no heap.", size=SMALL_SIZE, space_before=Pt(8))

    # =================================================================
    # Bounds-Checked CDR — page 2: align method
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "CDR Serialization \u2014 Alignment Check")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "Every write checks remaining capacity before touching the buffer:")
    add_code_block(tf, """\
pub fn align(&mut self, alignment: usize) -> Result<(), SerError> {
    let offset = self.pos - self.origin;
    let padding = (alignment - (offset % alignment)) % alignment;
    if self.remaining() < padding {
        return Err(SerError::BufferTooSmall); // Caught, not UB
    }
    for i in 0..padding {
        self.buf[self.pos + i] = 0;
    }
    self.pos += padding;
    Ok(())
}""", lang="rust")
    add_body_para(tf, "Buffer overflow is a compile-time-sized, runtime-checked error \u2014 never undefined behavior.", size=SMALL_SIZE, space_before=Pt(6))

    # =================================================================
    # Deterministic Executor Scheduling (page 1)
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Deterministic Executor Scheduling")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "The executor uses a fixed-size callback arena \u2014 no heap allocation during operation.")
    add_code_block(tf, """\
/// Type-erased callback metadata (no Box, no dyn Trait)
pub(crate) struct CallbackMeta {
    pub(crate) offset: usize,          // Position in arena
    pub(crate) kind: EntryKind,        // Sub / Srv / Timer / Action
    pub(crate) try_process:
        unsafe fn(*mut u8, u64) -> Result<bool, TransportError>,
    pub(crate) has_data:  unsafe fn(*const u8) -> bool,
    pub(crate) pre_sample: unsafe fn(*mut u8), // LET pre-sampling
    pub(crate) invocation: InvocationMode,
    pub(crate) drop_fn:    unsafe fn(*mut u8),
}""", lang="rust")

    # =================================================================
    # Deterministic Executor Scheduling (page 2)
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Deterministic Executor \u2014 Inline Entries")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "Each callback entry is stored inline in the arena:")
    add_code_block(tf, """\
#[repr(C)]
pub(crate) struct SubInfoEntry<M, F, const RX_BUF: usize> {
    pub(crate) handle: RmwSubscriber,
    pub(crate) buffer: [u8; RX_BUF],  // Static receive buffer
    pub(crate) sampled_len: usize,
    pub(crate) callback: F,
    pub(crate) _phantom: PhantomData<M>,
}""", lang="rust")
    add_body_para(tf, "All sizes known at compile time. Arena capacity set via environment variables (NROS_EXECUTOR_MAX_CBS=64), enforced at build time.", space_before=Pt(8))

    # =================================================================
    # RTIC Integration — page 1: setup
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "RTIC Integration \u2014 Hardware-Scheduled Tasks")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "On Cortex-M, nano-ros integrates with RTIC v2. No software scheduler \u2014 the NVIC enforces priorities.")
    add_code_block(tf, """\
#![no_std]
#![no_main]

#[rtic::app(device = stm32f4xx_hal::pac,
            dispatchers = [USART1, USART2])]
mod app {
    use nros::prelude::*;
    use std_msgs::msg::Int32;

    #[local]
    struct Local {
        executor: Executor,
        publisher: EmbeddedPublisher<Int32>,
    }""", lang="rust")
    add_body_para(tf, "All handles are #[local] \u2014 no locks required.", size=SMALL_SIZE, space_before=Pt(6))

    # =================================================================
    # RTIC Integration — page 2: async task
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "RTIC Integration \u2014 Async Task")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "RTIC async tasks run on hardware priorities \u2014 no OS scheduler:")
    add_code_block(tf, """\
    #[task(local = [executor], priority = 1)]
    async fn net_poll(cx: net_poll::Context) {
        loop {
            cx.local.executor.spin_once(0);
            Mono::delay(10.millis()).await;
        }
    }

    #[task(local = [publisher], priority = 1)]
    async fn publish(cx: publish::Context) {
        let mut counter: i32 = 0;
        loop {
            counter = counter.wrapping_add(1);
            cx.local.publisher
                .publish(&Int32 { data: counter }).ok();
            Mono::delay(1000.millis()).await;
        }
    }
}""", lang="rust")
    add_body_para(tf, "RTIC guarantees data-race freedom through ownership, not mutexes.", size=SMALL_SIZE, space_before=Pt(4))

    # =================================================================
    # E2E Safety Protocol
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "E2E Safety Protocol (EN 50159)")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "nano-ros treats the transport as an untrusted black channel.")

    add_table(
        slide,
        left=Inches(0.5), top=Inches(2.0), width=Inches(9.0),
        headers=["Threat", "Defense", "Implementation"],
        rows=[
            ["Corruption", "CRC-32", "Over CDR payload, in attachment"],
            ["Repetition", "Sequence number", "Subscriber tracks expected seq"],
            ["Deletion", "Seq gap detection", "Gap > 0 flagged immediately"],
            ["Insertion", "Seq + source GID", "Rejects unexpected seq jumps"],
            ["Resequencing", "Monotonic seq", "Non-monotonic = error"],
            ["Delay", "Timestamp", "Stale messages rejected"],
            ["Masquerade", "Publisher GID", "Random 16-byte identity"],
        ],
    )

    add_body_para(tf, "Enabled via a single feature flag:", space_before=Pt(6))
    add_code_block(tf, 'safety-e2e = ["nros-rmw/safety-e2e", "nros-rmw-zenoh?/safety-e2e"]', lang="toml")
    add_code_block(tf, """\
assert!(output.matches("crc=ok").count() >= 3);
assert_eq!(output.matches("crc=FAIL").count(), 0);
assert!(output.matches("seq_gap=0").count() >= 3);""", lang="rust")

    # =================================================================
    # TSN — Feasible Path (page 1)
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "TSN \u2014 Feasible Path")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "TSN requires hardware support. nano-ros\u2019s path: ThreadX + NetX Duo (TSN APIs included).")
    add_body_para(tf, "The ThreadX network driver is already implemented:", bold=True)
    add_code_block(tf, """\
// network.c — zenoh-pico transport over NetX Duo BSD sockets
#include "tx_api.h"
#include "nx_api.h"
#include "nxd_bsd.h"

z_result_t _z_open_tcp(_z_sys_net_socket_t *sock,
                       const _z_sys_net_endpoint_t rep, ...) {
    sock->_fd = nx_bsd_socket(AF_INET, SOCK_STREAM, IPPROTO_TCP);
    struct nx_bsd_sockaddr_in addr;
    _z_ep_to_sockaddr(&rep, &addr);
    int rc = nx_bsd_connect(sock->_fd,
        (struct nx_bsd_sockaddr *)&addr, sizeof(addr));
    return (rc < 0) ? _Z_ERR_GENERIC : _Z_RES_OK;
}""", lang="c")

    # =================================================================
    # TSN — Upgrade Path (page 2)
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "TSN \u2014 Upgrade Path")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "The upgrade path to TSN (NetX Duo already provides the APIs):")

    add_table(
        slide,
        left=Inches(0.5), top=Inches(2.0), width=Inches(9.0),
        headers=["Layer", "Status", "Details"],
        rows=[
            ["nano-ros app", "Working", "zenoh-pico pub/sub + RPC"],
            ["NetX Duo BSD sockets", "Working", "TCP + UDP unicast transport"],
            ["NetX Duo TSN module", "API available", "CBS (802.1Qav), TAS (802.1Qbv), PTP (802.1AS)"],
            ["TSN Ethernet MAC", "Hardware upgrade", "NXP i.MX RT1170/1180, Infineon AURIX TC3xx"],
        ],
        font_size=Pt(11),
    )

    add_body_para(tf, "The software stack is ready. The remaining step is hardware with a TSN-capable MAC.", bold=True, space_before=Pt(8))

    # =================================================================
    # Verus Proofs — page 1: OneShot
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Formal Verification \u2014 Verus Proofs")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "Unbounded deductive proofs \u2014 mathematical guarantees for all inputs. 67 proofs across scheduling, CDR, E2E.")
    add_code_block(tf, """\
verus! {

/// OneShot timer fires exactly once, then becomes inert.
proof fn timer_oneshot_fires_once(s: TimerGhost, delta_ms: u64)
    requires s.mode is OneShot,
    ensures
        timer_fire_mode(s) is Inert,
        timer_fire_elapsed(s) == 0,
        !timer_update_ready(timer_after_fire(s), delta_ms),
{}

}""", lang="rust")
    add_body_para(tf, "Safety-critical one-time actions (e.g., emergency stop) can never repeat.", size=SMALL_SIZE, space_before=Pt(6))

    # =================================================================
    # Verus Proofs — page 2: Drift-free repeating
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Verus Proofs \u2014 Drift-Free Scheduling")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "Repeating timer preserves overshoot \u2014 no cumulative drift:")
    add_code_block(tf, """\
verus! {

/// Control loops fire at t=0, P, 2P, 3P... not t=0, P+e, 2P+2e...
proof fn timer_repeating_drift_free(s: TimerGhost, delta_ms: u64)
    requires
        s.mode is Repeating, s.period_ms > 0,
        sat_add(s.elapsed_ms, delta_ms) >= s.period_ms,
    ensures
        timer_fire_elapsed(TimerGhost {
            elapsed_ms: sat_add(s.elapsed_ms, delta_ms), ..s
        }) == sat_sub(
            sat_add(s.elapsed_ms, delta_ms), s.period_ms),
{}

}""", lang="rust")
    add_body_para(tf, "Excess time carries over to the next period \u2014 control loops stay aligned.", size=SMALL_SIZE, space_before=Pt(6))

    # =================================================================
    # Kani Model Checking — page 1: harness
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Formal Verification \u2014 Kani Model Checking")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "Bounded model checking \u2014 exhaustive exploration of all execution paths. 82 harnesses across core crates.")
    add_code_block(tf, """\
#[cfg(kani)]
mod verification {
    use super::*;

    /// NaN velocity is treated as stopped (safe default).
    #[kani::proof]
    fn nan_velocity_treated_as_stopped() {
        let threshold: f64 = kani::any();
        kani::assume(threshold > 0.0 && threshold.is_finite());

        let filter = StopFilter::new(threshold, threshold, threshold);
        let result = filter.apply(&make_twist(f64::NAN, 0.0, 0.0));

        assert_eq!(result.linear.x, 0.0);  // NaN → stopped
        assert_eq!(result.linear.y, 0.0);
        assert_eq!(result.angular.z, 0.0);
    }
}""", lang="rust")

    # =================================================================
    # Kani — page 2: NaN safety pattern
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Kani \u2014 NaN Safety Pattern")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "NaN comparisons require careful handling in safety-critical code:")
    add_code_block(tf, """\
// WRONG: NaN.abs() < threshold → false → NaN passes through
linear.x.abs() < self.vx_threshold

// RIGHT: !(NaN.abs() >= threshold) → !(false) → true → NaN treated as stopped
!(linear.x.abs() >= self.vx_threshold)""", lang="rust")
    add_body_para(tf, "Kani proves this holds for every possible f64 \u2014 including NaN, Inf, and subnormals.", space_before=Pt(8))

    # =================================================================
    # Compile-Time Capacity
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Compile-Time Capacity via Environment Variables")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "All buffer sizes and arena capacities are set at compile time \u2014 no runtime allocation decisions.")
    add_code_block(tf, """\
NROS_EXECUTOR_MAX_CBS=64             # Max callback slots
NROS_MAX_PARAMETERS=64               # Max ROS 2 parameters
NROS_PARAM_SERVICE_BUFFER_SIZE=8192  # CDR buffer per param service

ZPICO_MAX_PUBLISHERS=40              # zenoh-pico publisher slots
ZPICO_MAX_SUBSCRIBERS=16             # zenoh-pico subscriber slots
ZPICO_MAX_LIVELINESS=64              # zenoh-pico liveliness tokens""", lang="bash")
    add_body_para(tf, "Read by build.rs and baked into the binary as constants. Exceeding a limit returns an error \u2014 not a heap allocation.", space_before=Pt(8))
    add_body_para(tf, "Autoware Sentinel: 37 publishers, 9 subscribers, 1 timer, 11 services, 62 parameters \u2014 zero heap.", bold=True, space_before=Pt(6))

    # =================================================================
    # Real-Time Lint Guide
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Real-Time Lint Guide")
    tf = get_body(slide)
    clear_body(tf)
    add_body_para(tf, "Clippy lints detect common real-time anti-patterns at build time:")
    add_code_block(tf, """\
cargo clippy -- \\
    -D clippy::infinite_iter \\
    -D clippy::while_immutable_condition \\
    -D clippy::never_loop \\
    -D clippy::empty_loop \\
    -W clippy::large_stack_arrays \\
    -W clippy::large_types_passed_by_value""", lang="bash")
    add_body_para(tf, "Edition 2024 safety improvements:", bold=True, space_before=Pt(8))
    add_bullet(tf, 'unsafe extern "C" { ... } blocks require the unsafe keyword')
    add_bullet(tf, "#[unsafe(no_mangle)] replaces #[no_mangle]")
    add_bullet(tf, "Explicit unsafe { ... } required inside unsafe fn bodies")
    add_body_para(tf, "The compiler is the first reviewer.", space_before=Pt(6))

    # =================================================================
    # Summary — Defense in Depth (diagram image)
    # =================================================================
    slide = new_slide(prs, LAYOUT_TITLE_BODY, "Summary \u2014 Defense in Depth")
    tf = get_body(slide)
    clear_body(tf)

    # Generate stack diagram image
    diagram = _make_stack_diagram([
        ("FORMAL VERIFICATION", "67 Verus proofs  \u00b7  82 Kani harnesses"),
        ("E2E SAFETY PROTOCOL  \u00b7  TSN", "CRC-32  \u00b7  sequence tracking  \u00b7  EN 50159  \u00b7  802.1Qbv"),
        ("DETERMINISTIC SCHEDULING", "Fixed arena  \u00b7  RTIC integration  \u00b7  LET semantics"),
        ("SAFE SERIALIZATION", "Bounds-checked CDR  \u00b7  no reallocation  \u00b7  Result<T,E>"),
        ("SAFE TYPES  (Rust + C++)", "heapless::Vec<T,N>  \u00b7  FixedSequence<T,N>  \u00b7  NROS_TRY"),
        ("COMPILER ENFORCEMENT", "#![no_std]  \u00b7  no alloc  \u00b7  feature gates  \u00b7  Edition 2024"),
        ("RUST LANGUAGE", "Ownership  \u00b7  borrowing  \u00b7  async state machines  \u00b7  Send+Sync"),
    ])
    slide.shapes.add_picture(diagram, Inches(0.5), Inches(1.3), Inches(9.0))

    add_body_para(tf, "Together: safety guarantees from language level through mathematical proof \u2014 with zero heap allocation.", bold=True, space_before=Pt(8))

    # =================================================================
    # Save
    # =================================================================
    prs.save(str(OUTPUT))
    print(f"Saved {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
